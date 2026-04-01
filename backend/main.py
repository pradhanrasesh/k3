# backend/main.py
# Unified backend entrypoint for COA Redaction Tool

import os
import json
import re
from typing import List, Optional

from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from fastapi import HTTPException

from backend.ocr_report import app as ocr_app
from backend.api.company_detection import router as company_router
from backend.api.auto_suggest import router as auto_suggest_router
from backend.api.routes.redaction_barcodes import router as barcode_router
from backend.api.ocr import router as ocr_router
from backend.api.document_review import router as document_review_router
from backend.template_loader import TemplateLoader
from backend.rules.merge_engine import detect_company
from backend.api.ai_training import router as ai_training_router
from backend.api.document_security import router as document_security_router
from backend.api.verification import router as verification_router
from backend.api.page_formatting import router as page_formatting_router
from backend.api.extraction import router as extraction_router
from backend.api.removal import router as removal_router
from backend.api.automation import router as automation_router
from backend.api.general import router as general_router
from backend.api.advanced_formatting import router as advanced_formatting_router
from backend.api.developer_tools import router as developer_tools_router
from backend.pdf_utils import PDFUtils
from backend.studio_rule_bridge import build_studio_rule_payload, merge_template_data, load_template_file
from backend.exceptions import (
    RedactionError,
    PDFProcessingError,
    OCRProcessingError,
    TemplateError,
    ValidationError,
    ConfigurationError,
    redaction_error_handler,
    http_exception_handler,
    generic_exception_handler,
)
import shutil
import uuid

app = FastAPI()

# ------------------------------------------------------------
# PDF TOOLS API
# ------------------------------------------------------------
@app.post("/api/tools/merge")
async def api_merge_pdfs(files: List[UploadFile] = File(...)):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    paths = []
    try:
        for f in files:
            path = os.path.join(temp_dir, f.filename)
            with open(path, "wb") as buffer:
                shutil.copyfileobj(f.file, buffer)
            paths.append(path)
        
        out_name = "merged.pdf"
        out_path = os.path.join(temp_dir, out_name)
        PDFUtils.merge_pdfs(paths, out_path)
        
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename=out_name, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/merge-advanced")
async def api_merge_pdfs_advanced(
    files: List[UploadFile] = File(...),
    page_ranges: Optional[str] = Form(None)
):
    """
    Merge multiple PDFs with optional page ranges per file.
    page_ranges: comma-separated list of page range strings, one per file.
    Example: "1-10, 14, 25-", "all", "1-5"
    If fewer ranges than files, missing ranges default to all pages.
    """
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    paths = []
    try:
        for f in files:
            path = os.path.join(temp_dir, f.filename)
            with open(path, "wb") as buffer:
                shutil.copyfileobj(f.file, buffer)
            paths.append(path)
        
        out_name = "merged.pdf"
        out_path = os.path.join(temp_dir, out_name)
        
        # Parse page ranges
        range_list = None
        if page_ranges:
            # Split by comma, but we need to handle nested commas within ranges.
            # Simple approach: split by ';' or treat as list of strings where each file range is separated by ';'
            # We'll use ';' as delimiter for simplicity.
            range_list = [r.strip() for r in page_ranges.split(';') if r.strip()]
        
        PDFUtils.merge_pdfs(paths, out_path, page_ranges=range_list)
        
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename=out_name, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/mix")
async def api_mix_pdfs(
    files: List[UploadFile] = File(...),
    reverse_order: bool = Form(False)
):
    """
    Mix PDFs by taking pages alternately from each input file.
    Perfect for single-sided scans.
    reverse_order: If True, take pages in reverse order (last page first).
    """
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    paths = []
    try:
        for f in files:
            path = os.path.join(temp_dir, f.filename)
            with open(path, "wb") as buffer:
                shutil.copyfileobj(f.file, buffer)
            paths.append(path)
        
        out_name = "mixed.pdf"
        out_path = os.path.join(temp_dir, out_name)
        PDFUtils.mix_pdfs(paths, out_path, reverse_order=reverse_order)
        
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename=out_name, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/split")
async def api_split_pdf(file: UploadFile = File(...)):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        paths = PDFUtils.split_pdf(path, os.path.join(temp_dir, "split"))
        # Return a ZIP of the split pages? 
        # For simplicity in this step, let's just return the count and paths
        return {"status": "ok", "pages": len(paths), "message": "Split complete. Zip download coming soon."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/rotate")
async def api_rotate_pdf(file: UploadFile = File(...), degrees: int = Form(90)):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        out_path = os.path.join(temp_dir, "rotated_" + file.filename)
        PDFUtils.rotate_pdf(path, out_path, degrees)
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="rotated_" + file.filename, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/compress")
async def api_compress_pdf(file: UploadFile = File(...)):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        out_path = os.path.join(temp_dir, "compressed_" + file.filename)
        PDFUtils.compress_pdf(path, out_path)
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="compressed_" + file.filename, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/password/add")
async def api_add_password(file: UploadFile = File(...), password: str = Form(...)):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        out_path = os.path.join(temp_dir, "locked_" + file.filename)
        PDFUtils.add_password(path, out_path, password)
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="locked_" + file.filename, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/convert/pdf-to-images")
async def api_pdf_to_images(file: UploadFile = File(...)):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        paths = PDFUtils.pdf_to_images(path, temp_dir)
        if paths:
            from fastapi.responses import FileResponse
            return FileResponse(paths[0], filename="page_1.png", media_type="image/png")
        return {"status": "ok", "message": "No pages found"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/convert/images-to-pdf")
async def api_images_to_pdf(files: List[UploadFile] = File(...)):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    paths = []
    try:
        for f in files:
            path = os.path.join(temp_dir, f.filename)
            with open(path, "wb") as buffer:
                shutil.copyfileobj(f.file, buffer)
            paths.append(path)
        
        out_path = os.path.join(temp_dir, "converted.pdf")
        PDFUtils.images_to_pdf(paths, out_path)
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="converted.pdf", media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/metadata")
async def api_update_metadata(file: UploadFile = File(...), title: str = Form(None), author: str = Form(None)):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        out_path = os.path.join(temp_dir, "meta_" + file.filename)
        PDFUtils.update_metadata(path, out_path, {"title": title, "author": author})
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="meta_" + file.filename, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/tools/watermark")
async def api_add_watermark(file: UploadFile = File(...), text: str = Form("WATERMARK")):
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        out_path = os.path.join(temp_dir, "watermarked_" + file.filename)
        PDFUtils.add_watermark(path, out_path, text)
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="watermarked_" + file.filename, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ------------------------------------------------------------
# Additional PDF Tools (ilovepdf.com features)
# ------------------------------------------------------------

@app.post("/api/tools/pdf-to-word")
async def api_pdf_to_word(file: UploadFile = File(...)):
    """Convert PDF to Word document (DOCX)."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Use the existing plugin
        from backend.plugins.ocr_pdf_to_word import OCRWordPlugin
        plugin = OCRWordPlugin()
        out_path = plugin.run(path, {})
        
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename=os.path.basename(out_path), media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to Word conversion failed: {str(e)}")

@app.post("/api/tools/pdf-to-pptx")
async def api_pdf_to_pptx(file: UploadFile = File(...)):
    """Convert PDF to PowerPoint presentation (PPTX)."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Simple implementation: extract images and create a presentation
        import fitz
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        
        doc = fitz.open(path)
        for i, page in enumerate(doc):
            slide = prs.slides.add_slide(blank_slide_layout)
            # Convert page to image
            pix = page.get_pixmap()
            img_path = os.path.join(temp_dir, f"page_{i+1}.png")
            pix.save(img_path)
            
            # Add image to slide
            left = Inches(0.5)
            top = Inches(0.5)
            slide.shapes.add_picture(img_path, left, top, height=Inches(7))
        
        out_path = os.path.join(temp_dir, "converted.pptx")
        prs.save(out_path)
        
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="converted.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except ImportError:
        raise HTTPException(status_code=501, detail="PPTX library not installed. Install 'python-pptx' package.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to PowerPoint conversion failed: {str(e)}")

@app.post("/api/tools/pdf-to-excel")
async def api_pdf_to_excel(file: UploadFile = File(...)):
    """Convert PDF to Excel spreadsheet (XLSX)."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Extract tables from PDF using camelot or tabula
        import fitz
        import pandas as pd
        
        doc = fitz.open(path)
        all_tables = []
        
        for i, page in enumerate(doc):
            # Try to extract text as table (simplified)
            text = page.get_text("text")
            if text.strip():
                # Simple approach: split lines and columns by tabs/whitespace
                lines = text.strip().split('\n')
                if lines:
                    # Create a simple DataFrame
                    df = pd.DataFrame([line.split('\t') for line in lines if line.strip()])
                    all_tables.append(df)
        
        if not all_tables:
            raise HTTPException(status_code=400, detail="No extractable table data found in PDF")
        
        # Combine tables or use the first one
        combined_df = pd.concat(all_tables, ignore_index=True)
        
        out_path = os.path.join(temp_dir, "converted.xlsx")
        combined_df.to_excel(out_path, index=False)
        
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="converted.xlsx", media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    except ImportError as e:
        raise HTTPException(status_code=501, detail=f"Required library not installed: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to Excel conversion failed: {str(e)}")

@app.post("/api/tools/word-to-pdf")
async def api_word_to_pdf(file: UploadFile = File(...)):
    """Convert Word document to PDF."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Use libreoffice or python-docx with reportlab
        # For now, return a placeholder
        raise HTTPException(status_code=501, detail="Word to PDF conversion not yet implemented. Requires libreoffice or reportlab.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Word to PDF conversion failed: {str(e)}")

@app.post("/api/tools/pptx-to-pdf")
async def api_pptx_to_pdf(file: UploadFile = File(...)):
    """Convert PowerPoint presentation to PDF."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Use libreoffice or python-pptx with reportlab
        raise HTTPException(status_code=501, detail="PowerPoint to PDF conversion not yet implemented. Requires libreoffice or reportlab.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PowerPoint to PDF conversion failed: {str(e)}")

@app.post("/api/tools/excel-to-pdf")
async def api_excel_to_pdf(file: UploadFile = File(...)):
    """Convert Excel spreadsheet to PDF."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Use pandas with reportlab or libreoffice
        raise HTTPException(status_code=501, detail="Excel to PDF conversion not yet implemented. Requires libreoffice or reportlab.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Excel to PDF conversion failed: {str(e)}")

@app.post("/api/tools/html-to-pdf")
async def api_html_to_pdf(url: str = Form(...), html: str = Form(None)):
    """Convert HTML/URL to PDF."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    
    try:
        # Use weasyprint or pdfkit
        raise HTTPException(status_code=501, detail="HTML to PDF conversion not yet implemented. Requires weasyprint or pdfkit.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"HTML to PDF conversion failed: {str(e)}")

@app.post("/api/tools/pdf-to-pdfa")
async def api_pdf_to_pdfa(file: UploadFile = File(...)):
    """Convert PDF to PDF/A format for archiving."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Use pikepdf or veraPDF
        import pikepdf
        from pikepdf import Pdf
        
        out_path = os.path.join(temp_dir, "pdfa_" + file.filename)
        with Pdf.open(path) as pdf:
            # Simple conversion - in reality need proper PDF/A validation
            pdf.save(out_path)
        
        from fastapi.responses import FileResponse
        return FileResponse(out_path, filename="pdfa_" + file.filename, media_type="application/pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to PDF/A conversion failed: {str(e)}")

@app.post("/api/tools/compare")
async def api_compare_pdfs(file1: UploadFile = File(...), file2: UploadFile = File(...)):
    """Compare two PDFs and highlight differences."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path1 = os.path.join(temp_dir, file1.filename)
    path2 = os.path.join(temp_dir, file2.filename)
    
    with open(path1, "wb") as buffer:
        shutil.copyfileobj(file1.file, buffer)
    with open(path2, "wb") as buffer:
        shutil.copyfileobj(file2.file, buffer)
    
    try:
        # Use difflib or specialized PDF comparison
        import fitz
        
        doc1 = fitz.open(path1)
        doc2 = fitz.open(path2)
        
        differences = []
        max_pages = max(len(doc1), len(doc2))
        
        for i in range(max_pages):
            if i < len(doc1) and i < len(doc2):
                text1 = doc1[i].get_text("text") if i < len(doc1) else ""
                text2 = doc2[i].get_text("text") if i < len(doc2) else ""
                if text1 != text2:
                    differences.append(f"Page {i+1}: Content differs")
            elif i >= len(doc1):
                differences.append(f"Page {i+1}: Only in second PDF")
            else:
                differences.append(f"Page {i+1}: Only in first PDF")
        
        doc1.close()
        doc2.close()
        
        return {
            "status": "ok",
            "differences": differences,
            "summary": f"Found {len(differences)} differences"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF comparison failed: {str(e)}")

@app.post("/api/tools/ai-summarize")
async def api_ai_summarize(file: UploadFile = File(...)):
    """Generate AI summary of PDF content."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Extract text
        import fitz
        doc = fitz.open(path)
        text = ""
        for page in doc:
            text += page.get_text("text") + "\n"
        doc.close()
        
        # Simple summary (first 500 chars) - in reality would use AI model
        summary = text[:500] + "..." if len(text) > 500 else text
        
        return {
            "status": "ok",
            "summary": summary,
            "original_length": len(text),
            "summary_length": len(summary)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI summarization failed: {str(e)}")

@app.post("/api/tools/translate")
async def api_translate_pdf(file: UploadFile = File(...), target_lang: str = Form("en")):
    """Translate PDF content to another language."""
    temp_dir = os.path.join(PROJECT_ROOT, "temp_tools", str(uuid.uuid4()))
    os.makedirs(temp_dir, exist_ok=True)
    path = os.path.join(temp_dir, file.filename)
    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    
    try:
        # Extract text
        import fitz
        doc = fitz.open(path)
        text = ""
        for page in doc:
            text += page.get_text("text") + "\n"
        doc.close()
        
        # Simple translation stub - in reality would use translation API
        translated = f"[Translated to {target_lang}]: {text[:200]}..."
        
        return {
            "status": "ok",
            "original_text_sample": text[:200],
            "translated_text_sample": translated,
            "target_language": target_lang
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF translation failed: {str(e)}")

# ------------------------------------------------------------
# CORS (allow frontend at 127.0.0.1:5500)
# ------------------------------------------------------------
origins = [
    "http://127.0.0.1:5500",
    "http://localhost:5500",
    "*",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ------------------------------------------------------------
# Exception handlers (inspired by Stirling-PDF)
# ------------------------------------------------------------
app.add_exception_handler(RedactionError, redaction_error_handler)
app.add_exception_handler(HTTPException, http_exception_handler)
app.add_exception_handler(Exception, generic_exception_handler)

# ------------------------------------------------------------
# Health check (used by dev + smoke-tests)
# ------------------------------------------------------------
@app.get("/health")
def health_check():
    return {"status": "ok", "service": "COA Redaction Tool"}

# ------------------------------------------------------------
# Include routers
# ------------------------------------------------------------
# Company detection router (/company/detect)
app.include_router(company_router)

# Unified rule-engine auto-suggest endpoint:
#   POST /redact/template
app.include_router(auto_suggest_router)

# Offline AI (local training) endpoint:
#   POST /api/ai/learn
app.include_router(ai_training_router)

# Pure barcode/QR endpoint:
#   POST /api/redact/auto-suggest-barcodes
app.include_router(barcode_router, prefix="/api")

# OCR engine endpoint:
#   POST /ocr
app.include_router(ocr_router, prefix="/api")

# Document Review endpoints:
#   POST /api/document/...
app.include_router(document_review_router)

# Document Security endpoints:
#   POST /api/security/...
app.include_router(document_security_router, prefix="/api/security")

# Verification endpoints:
#   POST /api/verification/...
app.include_router(verification_router, prefix="/api/verification")

# Page Formatting endpoints:
#   POST /api/formatting/...
app.include_router(page_formatting_router, prefix="/api/formatting")

# Extraction endpoints:
#   POST /api/extraction/...
app.include_router(extraction_router, prefix="/api/extraction")

# Removal endpoints:
#   POST /api/removal/...
app.include_router(removal_router, prefix="/api/removal")

# Automation endpoints:
#   POST /api/automation/...
app.include_router(automation_router, prefix="/api/automation")

# General PDF operations endpoints:
#   POST /api/general/...
app.include_router(general_router, prefix="/api/general")

# Advanced Formatting endpoints:
#   POST /api/advanced/...
app.include_router(advanced_formatting_router, prefix="/api/advanced")

# Developer Tools endpoints:
#   POST /api/developer/...
app.include_router(developer_tools_router, prefix="/api/developer")

# ------------------------------------------------------------
# Fallback /detect-company endpoint (used by Template_Detect_Backend.js)
# ------------------------------------------------------------
PROJECT_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
COMPANY_RULES_DIR = os.path.join(PROJECT_ROOT, "config", "rules", "company_rules")

loader = TemplateLoader(templates_dir=COMPANY_RULES_DIR)


def _build_company_template_raw_from_saved_rule(saved: dict) -> dict:
    """
    Convert frontend SaveRule.js payload into the company JSON schema
    that backend/rules/merge_engine expects.
    """
    company_id = (saved.get("company_id") or "").strip()
    display_name = (saved.get("display_name") or company_id or "Unnamed Company").strip()

    rules = saved.get("rules")
    if not isinstance(rules, list):
        rules = []

    regex_rules = []
    layout_rules = []
    sensitive_terms = []
    for idx, r in enumerate(rules):
        if not isinstance(r, dict):
            continue

        sample_text = (r.get("sample_text") or "").strip()
        if sample_text:
            sensitive_terms.append(sample_text)
        if not sample_text:
            sample_text = ""

        rid = str(r.get("id") or r.get("label") or f"saved_rule_{idx}").strip()
        if sample_text:
            pattern = re.escape(sample_text)
            regex_rules.append(
                {
                    "id": rid,
                    "pattern": pattern,
                    "action": "suggest",
                }
            )

        rects = r.get("rects") or []
        if rects and isinstance(rects[0], dict):
            rect = rects[0]
            layout_rules.append(
                {
                    "id": f"{rid}_layout",
                    "label": r.get("label") or "Learned Area",
                    "rect": {
                        "x0": float(rect.get("x0", 0.0) or 0.0),
                        "y0": float(rect.get("y0", 0.0) or 0.0),
                        "x1": float(rect.get("x1", 1.0) or 1.0),
                        "y1": float(rect.get("y1", 1.0) or 1.0),
                    },
                    "relative": True,
                    "page_scope": "all",
                    "action": "suggest",
                }
            )

    payload = build_studio_rule_payload(company_id, display_name, rules)
    payload["regex"] = payload.get("regex", []) + regex_rules
    payload["layout"] = payload.get("layout", []) + layout_rules
    payload["sensitive_terms"] = list(dict.fromkeys((payload.get("sensitive_terms") or []) + sensitive_terms))
    return payload


@app.get("/api/templates")
def api_list_templates():
    templates_meta = []
    for company_id in loader.list_templates():
        tmpl = loader.get_template(company_id) or {}
        templates_meta.append(
            {
                "company_id": company_id,
                "display_name": tmpl.get("display_name", company_id),
            }
        )
    return {"templates": templates_meta}


@app.get("/api/templates/{company_id}")
def api_get_template(company_id: str):
    template = loader.get_template(company_id)
    if not template:
        raise HTTPException(status_code=404, detail="Template not found")
    return template


@app.post("/api/templates/save-rule")
async def api_save_rule(rule: str = Form(...)):
    """
    Saves a "saved rule" created from frontend suggestions.
    Frontend posts FormData field named `rule` as a JSON string.
    """
    try:
        saved = json.loads(rule)
    except json.JSONDecodeError:
        return JSONResponse({"error": "Invalid JSON in rule"}, status_code=400)

    company_id = (saved.get("company_id") or "").strip()
    if not company_id:
        return JSONResponse({"error": "Missing company_id"}, status_code=400)

    raw_template = _build_company_template_raw_from_saved_rule(saved)

    # Safety: ensure company_id is always set correctly.
    raw_template["company_id"] = company_id
    raw_template["display_name"] = raw_template.get("display_name") or company_id

    os.makedirs(COMPANY_RULES_DIR, exist_ok=True)
    path = os.path.join(COMPANY_RULES_DIR, f"{company_id}.json")

    merged_template = merge_template_data(load_template_file(path), raw_template)

    with open(path, "w", encoding="utf-8") as f:
        json.dump(merged_template, f, indent=2)

    # Reload so the UI immediately sees the newly saved template.
    loader.load_templates()
    return {
        "status": "ok",
        "saved_to": path,
        "company_id": company_id,
        "studio_rule": merged_template.get("studio_rule", {}),
    }


@app.post("/detect-company")
async def detect_company_endpoint(file: UploadFile = File(...)):
    """
    Simple backend company detection used by Template_Detect_Backend.js.
    - Extracts text from PDF with PyMuPDF (via api_server's PDFTextExtractor if needed)
    - Uses merge_engine.detect_company over config/rules/company_rules/*.json
    - Returns { company_id, display_name } or nulls
    """
    import fitz  # local import to avoid unused at module level

    pdf_bytes = await file.read()

    # Extract text from PDF
    text_chunks = []
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        for page in doc:
            text_chunks.append(page.get_text("text") or "")
        doc.close()
    except Exception as e:
        print("[detect-company] ERROR extracting text:", e)
        return {"company_id": None, "display_name": None}

    full_text = "\n".join(text_chunks)

    rules = detect_company(full_text, COMPANY_RULES_DIR)
    if not rules:
        return {"company_id": None, "display_name": None}

    company_id = rules.get("company_id") or None
    display_name = rules.get("display_name") or company_id

    return {
        "company_id": company_id,
        "display_name": display_name,
    }


# Mount OCR + Redaction + Plugin API LAST.
# Mounting at "/" can otherwise shadow non-mounted routes (notably /api/templates)
# depending on Starlette route ordering.
app.mount("/", ocr_app)
